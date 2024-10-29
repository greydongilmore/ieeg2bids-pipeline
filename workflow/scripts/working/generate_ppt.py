#!/usr/bin/env python3
# -*- coding: utf-8 -*-
import os
import glob

import pandas as pd
from pptx.enum.text import MSO_AUTO_SIZE
import numpy as np
import datetime


def add_slide(presentation, layout, title_dict):
    slide = presentation.slides.add_slide(layout)  # adding a slide
    for ititle in list(title_dict):
        title = slide.shapes.add_textbox(*title_dict[ititle]['position'])

        tf = title.text_frame
        tf.auto_size = MSO_AUTO_SIZE.NONE
        tf.word_wrap = False

        p = tf.paragraphs[0]
        p.alignment = PP_ALIGN.CENTER
        p.text = ititle
        p.font.size = Pt(title_dict[ititle]['font_size'])
        p.font.color.rgb = title_dict[ititle]['color']

    return slide


def write_cell(table, row, col, value):
    table.cell(row, col).text = "%s" % value


def format_table_header(tbl_tmp):

    col_cnt = 0
    header_vals = [
        [(0, col_cnt), (1, 0), 'Electrode']
    ]

    col_cnt += 1
    if len(tbl_tmp.table.rows[1].cells) == 8:
        header_vals.extend([[(0, col_cnt), (1, 1), 'Implanter']])
        col_cnt += 1

    header_vals.extend([
        [(0, col_cnt), (0, col_cnt+1), 'Target Error'],
        [(0, col_cnt+2), (0, col_cnt+3), 'Entry Error'],
        [(0, col_cnt+4), (1, col_cnt+4), 'Radial Angle'],
        [(0, col_cnt+5), (1, col_cnt+5), 'Line Angle'],
        [(1, col_cnt), 'Euclidean'],
        [(1, col_cnt+1), 'Radial'],
        [(1, col_cnt+2), 'Euclidean'],
        [(1, col_cnt+3), 'Radial'],
    ])

    for ihead in header_vals:
        if len(ihead) > 2:
            cell = tbl_tmp.table.cell(ihead[0][0], ihead[0][1])
            cell.merge(tbl_tmp.table.cell(ihead[1][0], ihead[1][1]))
            cell.text = ihead[2]
        else:
            cell = tbl_tmp.table.cell(ihead[0][0], ihead[0][1])
            cell.text = ihead[1]

        cell.vertical_anchor = MSO_ANCHOR.BOTTOM
        cell.text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
        cell.text_frame.paragraphs[0].font.bold = True
        cell.text_frame.paragraphs[0].font.color.rgb = RGBColor(255, 255, 255)
        cell.fill.solid()
        cell.fill.fore_color.rgb = RGBColor(51, 51, 51)

    return tbl_tmp


color_map = {
    "gray": (102, 102, 102),
    "grey": (102, 102, 102),
    "black": (0, 0, 0),
    "red": (255, 0, 0),
    "blue": (42, 96, 153),
    "purple": (128, 0, 128),
    "orange": (255, 128, 0),
    "yellow": (255, 255, 0),
    "brown": (43, 34, 2),
    "green": (0, 169, 51),
    "white": (255, 255, 255),
}

aborted_lang = {
    'skipped',
    'aborted'
}

remap_dict = {
    'Electrode label ("aborted" if skipped)': 'Electrode label',
    'Label (6 chr)': 'Label'
}


# %%


debug = True

if debug:
    class dotdict(dict):
        """dot.notation access to dictionary attributes"""
        __getattr__ = dict.get
        __setattr__ = dict.__setitem__
        __delattr__ = dict.__delitem__

    class Namespace:
        def __init__(self, **kwargs):
            self.__dict__.update(kwargs)

    isub = 'sub-P157'
    # data_dir = r'/media/greydon/lhsc_data/datasets/SEEG_peds/derivatives'
    data_dir = r'/home/greydon/Documents/data/SEEG/derivatives'

    input = dotdict({
                    'shopping_list': f'{data_dir}/seeg_scenes/{isub}/*shopping_list.xlsx',
                    'error_metrics': f'{data_dir}/seeg_scenes/{isub}/{isub}_error_metrics.xlsx',
                    })

    snakemake = Namespace(input=input)


pt_pin = 'PIN'
sx_date = 'yyyy-mm-dd'
lastname = "lastname"
firstname = "firstname"


if glob.glob(snakemake.input.shopping_list):
    df_elec_raw = pd.read_excel(
        glob.glob(snakemake.input.shopping_list)[0], header=None)
    df_elec = df_elec_raw.iloc[4:, :].reset_index(drop=True)

    # need to update the column names
    updated_colnames = df_elec_raw.iloc[3].values
    for idx, ilabel in [(i, x) for i, x in enumerate(updated_colnames) if x in list(remap_dict)]:
        updated_colnames[idx] = remap_dict[ilabel]

    df_elec.columns = updated_colnames
    if 'Electrode label' in list(df_elec):
        df_elec = df_elec[df_elec['Electrode label'] != 'aborted']
    elif 'Serial Num.' in list(df_elec):
        df_elec = df_elec[df_elec['Serial Num.'] != 'aborted']

    df_elec = df_elec.iloc[0:df_elec.loc[:, 'Target'].isnull().idxmax()]
    df_elec = df_elec[~df_elec['No.'].isnull()]
    df_elec = df_elec[~df_elec['Target'].isnull()]
    df_elec = df_elec.dropna(axis=1, how='all')

    if any(pd.isna(x) for x in list(df_elec)):
        df_elec.drop(np.nan, axis=1, inplace=True)

    if 'Ord.' in list(df_elec):
        if all(~df_elec.loc[:, 'Ord.'].isnull()):
            df_elec = df_elec.sort_values(by=['Ord.']).reset_index(drop=True)

    pin_idx = [i for i, x in enumerate(
        df_elec_raw.iloc[1].values) if x == 'PIN']
    if pin_idx:
        pt_pin = df_elec_raw.iloc[1, pin_idx[0]+1]

    sx_idx = [i for i, x in enumerate(
        df_elec_raw.iloc[2].values) if x == 'Date']
    if sx_idx:
        if isinstance(df_elec_raw.iloc[2, sx_idx[0]+1], datetime.datetime):
            sx_date = df_elec_raw.iloc[2, sx_idx[0]+1].strftime('%Y-%m-%d')
        elif '_' in df_elec_raw.iloc[2, sx_idx[0]+1]:
            sx_date = datetime.datetime.strptime(
                df_elec_raw.iloc[2, sx_idx[0]+1], '%Y_%m_%d').strftime('%Y-%m-%d')
        else:
            sx_date = datetime.datetime.strptime(
                df_elec_raw.iloc[2, sx_idx[0]+1], '%d/%b/%y').strftime('%Y-%m-%d')

    name_idx = [i for i, x in enumerate(
        df_elec_raw.iloc[0].values) if x == 'Name']
    if name_idx:
        if ',' in df_elec_raw.iloc[0, name_idx[0]+1]:
            lastname, firstname = df_elec_raw.iloc[0, name_idx[0]+1].split(',')
        else:
            firstname, lastname = df_elec_raw.iloc[0, name_idx[0]+1].split(' ')

        firstname = firstname.strip()
        lastname = lastname.strip()

elif os.path.exists(snakemake.input.error_metrics):
    df_elec_raw = pd.read_excel(snakemake.input.error_metrics, header=0)
    df_elec_raw = df_elec_raw.rename(columns={'electrode': 'Electrode label'})
    df_elec = df_elec_raw


prs = Presentation()
prs.slide_width = Inches(16)
prs.slide_height = Inches(9)

blank_slide_layout = prs.slide_layouts[6]
fill = blank_slide_layout.background.fill
fill.solid()
fill.fore_color.rgb = RGBColor(0, 0, 0)

# Title slide
title_dict = {
    f"{lastname}, {firstname}": {
        "font_size": 52,
        "color": RGBColor(255, 255, 255),
        "position": (Inches(3), Inches(2), Inches(10), Inches(1))
    },
    f"{pt_pin}": {
        "font_size": 36,
        "color": RGBColor(255, 255, 255),
        "position": (Inches(5), Inches(3), Inches(6), Inches(.8))
    },
    f"Implantation Date:\n{sx_date}": {
        "font_size": 36,
        "color": RGBColor(255, 255, 255),
        "position": (Inches(5), Inches(4.5), Inches(6), Inches(1.5))
    }
}

title_slide = add_slide(prs, blank_slide_layout, title_dict)
title_slide.name = "title slide"

if glob.glob(snakemake.input.shopping_list):
    # Shopping list
    shopping_list_slide = add_slide(prs, blank_slide_layout, {})
    shopping_list_slide.name = "shopping list"


# Errors
title_dict = {
    "Errors": {
        "font_size": 48,
        "color": RGBColor(255, 255, 255),
        "position": (Inches(3), Inches(.5), Inches(10), Inches(1))
    }
}

errors_data = None
if os.path.exists(snakemake.input.error_metrics):
    errors_data = pd.read_excel(snakemake.input.error_metrics, header=0)
    error_slide = add_slide(prs, prs.slide_layouts[6], title_dict)
    error_slide.name = "errors"
    error_slide.background.fill.solid()
    error_slide.background.fill.fore_color.rgb = RGBColor(0, 0, 0)

    width = Inches(13.0)
    height = Inches(5.0)
    left = (prs.slide_width - width) / 2
    top = (prs.slide_height - height) / 2

    tbl_tmp = error_slide.shapes.add_table(
        errors_data.shape[0]+2, errors_data.shape[1], left, top, width, height)
    tbl = format_table_header(tbl_tmp)

    for row in range(2, len(tbl.table.rows)):
        for cell in range(len(tbl.table.rows[row].cells)):
            if isinstance(errors_data.iloc[row-2, cell], str):
                tbl.table.rows[row].cells[cell].text_frame.text = errors_data.iloc[row-2, cell]
                tbl.table.rows[row].cells[cell].vertical_anchor = MSO_ANCHOR.BOTTOM
                tbl.table.rows[row].cells[cell].text_frame.paragraphs[0].font.bold = True
                tbl.table.rows[row].cells[cell].text_frame.paragraphs[0].font.color.rgb = RGBColor(
                    255, 255, 255)
                tbl.table.rows[row].cells[cell].fill.solid()
                tbl.table.rows[row].cells[cell].fill.fore_color.rgb = RGBColor(
                    51, 51, 51)
            elif isinstance(errors_data.iloc[row-2, cell], float):
                tbl.table.rows[row].cells[cell].text_frame.text = f"{errors_data.iloc[row-2,cell]:1.2f}"

                if errors_data.iloc[row-2, cell] <= 2:
                    col = RGBColor(99, 248, 99)
                elif errors_data.iloc[row-2, cell] > 2 and errors_data.iloc[row-2, cell] < 3:
                    col = RGBColor(255, 255, 0)
                else:
                    col = RGBColor(255, 95, 54)

                tbl.table.rows[row].cells[cell].fill.solid()
                tbl.table.rows[row].cells[cell].fill.fore_color.rgb = col
                tbl.table.rows[row].cells[cell].fill.fore_color.brightness = 0.4
                tbl.table.rows[row].cells[cell].text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
                tbl.table.rows[row].cells[cell].vertical_anchor = MSO_ANCHOR.BOTTOM

    tbl.left = int((prs.slide_width / 2) - (tbl.width / 2))
    tbl.top = int((prs.slide_height / 2) - (tbl.height / 2))


for _, row_elec in df_elec.iterrows():

    elec_label = None
    if any('Electrode label' in x for x in list(row_elec.keys())):
        elec_label = list(row_elec.keys())[[i for i,x in enumerate(list(row_elec.keys())) if 'Electrode label' in x][0]]
    elif 'Serial Num.' in list(row_elec.keys()):
        elec_label = 'Serial Num.'

    if not any(x in row_elec[list(row_elec.keys())] for x in aborted_lang):

        if any([x.lower() == 'label' for x in list(row_elec.keys())]):
            slide_title = f"{row_elec['Target']} ({row_elec['Label']})"
        elif any([x.lower() == 'target' for x in list(row_elec.keys())]):
            slide_title = row_elec['Target']
        else:
            slide_title = row_elec['Electrode label']

        title_dict = {
            slide_title: {
                "font_size": 48,
                "color": RGBColor(255, 255, 255),
                "position": (Inches(3), Inches(.5), Inches(10), Inches(1))
            }
        }

        elec_slide = add_slide(prs, blank_slide_layout, title_dict)
        elec_slide.name = slide_title

        if errors_data is not None:
            error_idx = []
            if any([x.lower() == 'label' for x in list(row_elec.keys())]):
                error_idx = [i for i, x in enumerate(list(
                    errors_data['electrode'].values)) if x.lower() in row_elec['Label'].lower()][0]
            elif any([x.lower() == 'label' for x in list(row_elec.keys())]):
                if [i for i, x in enumerate(errors_data['electrode']) if f'({x.lower()})' in row_elec['Target'].lower()]:
                    error_idx = [i for i, x in enumerate(
                        errors_data['electrode']) if f'({x.lower()})' in row_elec['Target'].lower()][0]
                elif [i for i, x in enumerate(errors_data['electrode']) if row_elec['Target'].lower().startswith(f'{x.lower()}')]:
                    error_idx = [i for i, x in enumerate(
                        errors_data['electrode']) if row_elec['Target'].lower().startswith(f'{x.lower()}')][0]
            else:
                if [i for i, x in enumerate(errors_data['electrode']) if f'({x.lower()})' in row_elec['Electrode label'].lower()]:
                    error_idx = [i for i, x in enumerate(
                        errors_data['electrode']) if f'({x.lower()})' in row_elec['Electrode label'].lower()][0]
                elif [i for i, x in enumerate(errors_data['electrode']) if row_elec['Electrode label'].lower().startswith(f'{x.lower()}')]:
                    error_idx = [i for i, x in enumerate(
                        errors_data['electrode']) if row_elec['Electrode label'].lower().startswith(f'{x.lower()}')][0]

            if isinstance(error_idx, int):

                width = Inches(13.0)
                height = Inches(1.25)
                left = (prs.slide_width - width) / 2
                top = (prs.slide_height - height) / 10

                tbl = elec_slide.shapes.add_table(
                    3, errors_data.shape[1], left, (top*9.5), width, height)
                tbl = format_table_header(tbl)

                cell = 0
                for ival in list(errors_data):
                    if isinstance(errors_data.loc[error_idx, ival], str):
                        tbl.table.rows[2].cells[cell].text_frame.text = errors_data.loc[error_idx, ival]
                        tbl.table.rows[2].cells[cell].vertical_anchor = MSO_ANCHOR.BOTTOM
                        tbl.table.rows[2].cells[cell].text_frame.paragraphs[0].font.bold = True
                        tbl.table.rows[2].cells[cell].text_frame.paragraphs[0].font.color.rgb = RGBColor(
                            255, 255, 255)
                        tbl.table.rows[2].cells[cell].fill.solid()
                        tbl.table.rows[2].cells[cell].fill.fore_color.rgb = RGBColor(
                            51, 51, 51)
                    elif isinstance(errors_data.loc[error_idx, ival], float):
                        tbl.table.rows[2].cells[cell].text_frame.text = f"{errors_data.loc[error_idx,ival]:1.2f}"

                        if errors_data.loc[error_idx, ival] <= 2:
                            col = RGBColor(99, 248, 99)
                        elif errors_data.loc[error_idx, ival] > 2 and errors_data.loc[error_idx, ival] < 3:
                            col = RGBColor(255, 255, 0)
                        else:
                            col = RGBColor(255, 95, 54)

                        tbl.table.rows[2].cells[cell].fill.solid()
                        tbl.table.rows[2].cells[cell].fill.fore_color.rgb = col
                        tbl.table.rows[2].cells[cell].fill.fore_color.brightness = 0.4
                        tbl.table.rows[2].cells[cell].text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
                        tbl.table.rows[2].cells[cell].vertical_anchor = MSO_ANCHOR.BOTTOM

                    cell += 1

        if isinstance(row_elec[elec_label], int):
            elec_color = 'black'
            elec_text = f"{row_elec[elec_label]}".zfill(3)
        elif '-' in row_elec[elec_label]:
            if row_elec[elec_label].split('-')[0].isdigit():
                elec_color = 'black'
                elec_text = f"{row_elec[elec_label]}".zfill(3)
        elif elec_label == 'Electrode label':
            elec_color = 'black'
            elec_text = f"{row_elec[elec_label]}".zfill(3)
        else:
            elec_color = ''.join(
                [x for x in row_elec[elec_label] if x.isalpha()]).lower()
            elec_text = row_elec[elec_label]

        textbox = elec_slide.shapes.add_textbox(
            Inches(13.5), Inches(4.5), Inches(2), Inches(.5))
        textbox.fill.solid()
        if elec_color in ("yellow", "green", "white"):
            textbox.fill.fore_color.rgb = RGBColor(0, 0, 0)
        else:
            textbox.fill.fore_color.rgb = RGBColor(255, 255, 255)
        tf = textbox.text_frame
        tf.auto_size = MSO_AUTO_SIZE.NONE
        tf.word_wrap = False

        p = tf.paragraphs[0]
        p.alignment = PP_ALIGN.CENTER
        p.text = elec_text
        p.font.size = Pt(24)
        p.font.bold = True
        p.font.color.rgb = RGBColor(
            color_map[elec_color][0], color_map[elec_color][1], color_map[elec_color][2])

        line = textbox.line
        line.color.rgb = RGBColor(255, 0, 0)
        line.width = Inches(0.04)

out_fname = f"{lastname.replace(' ','')}_{firstname}_{sx_date}_maps.pptx"
prs.save(f'{data_dir}/seeg_scenes/{isub}/{out_fname}')
